"""Services : Whisper via Groq (transcription) + extraction PPTX/PDF."""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Optional

from django.conf import settings

logger = logging.getLogger('notation')


# ============================================================================
# Whisper via Groq — transcription audio
# ============================================================================

def transcrire_audio(chemin_audio: str | Path, langue: str = 'fr') -> str:
    """Transcrit un fichier audio via Whisper hébergé sur Groq. Retourne le texte."""
    try:
        from groq import Groq
    except ImportError:
        logger.warning("groq package non installé")
        return ''

    if not settings.GROQ_API_KEY:
        logger.warning("GROQ_API_KEY non configurée, transcription désactivée")
        return ''

    client = Groq(api_key=settings.GROQ_API_KEY)
    # Mapping ISO langue
    langue_iso = {'fr': 'fr', 'en': 'en', 'ar': 'ar', 'es': 'es', 'de': 'de'}.get(langue, 'fr')
    chemin = Path(chemin_audio)
    try:
        with chemin.open('rb') as f:
            response = client.audio.transcriptions.create(
                model=settings.WHISPER_MODEL,
                file=f,
                language=langue_iso,
                response_format='text',
            )
        return response if isinstance(response, str) else getattr(response, 'text', '')
    except Exception:
        logger.exception("Erreur Whisper (Groq)")
        return ''


# ============================================================================
# Extraction de slides PPTX
# ============================================================================

def extraire_texte_pptx(chemin_pptx: str | Path) -> str:
    """Extrait tout le texte des slides PPTX (titres + contenu)."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx non installé")
        return ''

    chemin = Path(chemin_pptx)
    if not chemin.exists():
        return ''
    try:
        prs = Presentation(str(chemin))
        morceaux = []
        for i, slide in enumerate(prs.slides, 1):
            morceaux.append(f'\n--- SLIDE {i} ---')
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = ''.join(run.text for run in para.runs).strip()
                        if txt:
                            morceaux.append(txt)
        return '\n'.join(morceaux)
    except Exception:
        logger.exception("Erreur extraction PPTX")
        return ''


def _convertir_via_powerpoint_com(chemin: Path, dossier_sortie: Path) -> Optional[Path]:
    """Fallback Windows : conversion PPTX->PDF via PowerPoint COM (pywin32).
    Utilise uniquement si LibreOffice est absent."""
    try:
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()
        # msoFalse=0 pour WithWindow -> ouvre sans fenetre visible
        ppt_app = win32com.client.Dispatch('PowerPoint.Application')
        pdf_path = dossier_sortie / (chemin.stem + '.pdf')
        # WithWindow=False : ouverture invisible (requiert que l'app soit deja initialisee)
        presentation = ppt_app.Presentations.Open(
            str(chemin.resolve()), ReadOnly=True, Untitled=False, WithWindow=False
        )
        presentation.SaveAs(str(pdf_path.resolve()), 32)  # 32 = ppSaveAsPDF
        presentation.Close()
        ppt_app.Quit()
        pythoncom.CoUninitialize()
        if pdf_path.exists():
            logger.info("Conversion PowerPoint COM reussie : %s", pdf_path)
            return pdf_path
        logger.error("PDF PowerPoint introuvable apres conversion : %s", pdf_path)
        return None
    except ImportError:
        logger.warning("pywin32 non installe : fallback PowerPoint COM impossible")
        return None
    except Exception:
        logger.exception("Erreur conversion PowerPoint COM")
        return None


def convertir_pptx_en_pdf(chemin_pptx: str | Path, dossier_sortie: Path) -> Optional[Path]:
    """Convertit un fichier PPTX en PDF.
    Essaie dans l'ordre : LibreOffice headless, puis PowerPoint COM (Windows).
    Retourne le chemin du PDF produit, ou None en cas d'echec."""
    import subprocess
    import shutil

    chemin = Path(chemin_pptx)
    if not chemin.exists():
        logger.warning("PPTX introuvable : %s", chemin)
        return None

    dossier_sortie.mkdir(parents=True, exist_ok=True)

    # -- Tentative 1 : LibreOffice headless (Linux/Mac/Windows) ---------------
    candidates = ['libreoffice', 'soffice', 'libreoffice7.6', 'libreoffice7.5']
    exe = None
    for c in candidates:
        if shutil.which(c):
            exe = c
            break
    if exe is None:
        import os
        win_paths = [
            r'C:\Program Files\LibreOffice\program\soffice.exe',
            r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
        ]
        for wp in win_paths:
            if os.path.isfile(wp):
                exe = wp
                break

    if exe is not None:
        try:
            result = subprocess.run(
                [exe, '--headless', '--convert-to', 'pdf', '--outdir', str(dossier_sortie), str(chemin)],
                capture_output=True, text=True, timeout=120,
            )
            if result.returncode == 0:
                pdf_path = dossier_sortie / (chemin.stem + '.pdf')
                if pdf_path.exists():
                    return pdf_path
            logger.error("LibreOffice erreur (code %s) : %s", result.returncode, result.stderr)
        except subprocess.TimeoutExpired:
            logger.error("LibreOffice timeout sur %s", chemin)
        except Exception:
            logger.exception("Erreur appel LibreOffice")

    # -- Tentative 2 : PowerPoint COM (Windows uniquement) --------------------
    logger.info("LibreOffice indisponible, tentative via PowerPoint COM...")
    return _convertir_via_powerpoint_com(chemin, dossier_sortie)


# Alias conservé pour compatibilité avec les anciens appels (retourne liste vide si utilisé)
def slides_pptx_en_images(chemin_pptx: str | Path, dossier_sortie: Path) -> list[Path]:
    """Dépréciée — utiliser convertir_pptx_en_pdf(). Garde la signature pour compatibilité."""
    convertir_pptx_en_pdf(chemin_pptx, dossier_sortie)
    return []


# ============================================================================
# Extraction de PDF
# ============================================================================

def extraire_texte_pdf(chemin_pdf: str | Path) -> str:
    """Extrait le texte d'un PDF via PyMuPDF."""
    try:
        import fitz
    except ImportError:
        logger.warning("PyMuPDF non installé")
        return ''

    chemin = Path(chemin_pdf)
    if not chemin.exists():
        return ''
    try:
        doc = fitz.open(str(chemin))
        morceaux = []
        for i, page in enumerate(doc, 1):
            morceaux.append(f'\n--- PAGE {i} ---')
            morceaux.append(page.get_text())
        doc.close()
        return '\n'.join(morceaux)
    except Exception:
        logger.exception("Erreur extraction PDF")
        return ''


# ============================================================================
# Résumé des données de posture MediaPipe
# ============================================================================

def resumer_posture(donnees: list[dict]) -> dict:
    """Calcule des stats à partir d'une série de mesures face-api.js.

    Chaque entrée attendue :
      {timestamp, looking_at_camera, head_tilt, shoulders_visible, mouvement,
       expressions: {neutral, happy, sad, fearful, surprised, angry, disgusted}}

    Les champs expressions sont renseignés par faceExpressionNet (face-api.js 4s).
    Si absents (ancienne donnée), seules les stats de posture sont calculées.
    """
    if not donnees:
        return {}

    total = len(donnees)
    contact_visuel = sum(1 for d in donnees if d.get('looking_at_camera'))
    tilt_moyen = sum(d.get('head_tilt', 0) for d in donnees) / total

    # ── Agréger les expressions (moyennes sur toutes les mesures avec expressions) ──
    expr_keys = ['neutral', 'happy', 'sad', 'fearful', 'surprised', 'angry', 'disgusted']
    expr_totaux = {k: 0.0 for k in expr_keys}
    nb_avec_expr = 0
    for d in donnees:
        exprs = d.get('expressions', {})
        if exprs:
            nb_avec_expr += 1
            for k in expr_keys:
                expr_totaux[k] += exprs.get(k, 0)

    expressions_moyennes: dict = {}
    emotion_dominante = None
    if nb_avec_expr > 0:
        expressions_moyennes = {
            k: round(expr_totaux[k] / nb_avec_expr * 100, 1)
            for k in expr_keys
        }
        emotion_dominante = max(expressions_moyennes, key=expressions_moyennes.get)

    resume = {
        'nb_mesures': total,
        'pourcentage_contact_visuel': round(100 * contact_visuel / total, 1),
        'inclinaison_tete_moyenne_deg': round(tilt_moyen, 2),
        'commentaire_auto': _commentaire_posture(contact_visuel / total),
    }
    if expressions_moyennes:
        resume['expressions_moyennes_pct'] = expressions_moyennes
        resume['emotion_dominante'] = emotion_dominante
    return resume


def enrichir_posture_avec_prosody(resume: dict, prosody_raw: dict) -> dict:
    """Fusionne les données prosodiques (volume micro, ratio silence, variation) dans le résumé.

    prosody_raw attendu : {db_moyen: float, db_variation: float, ratio_silence: float}
    Ces données viennent de la Web Audio API (AnalyserNode côté navigateur) :
    - db_moyen : volume moyen (dB) — proxy de la puissance vocale
    - db_variation : écart-type du volume — proxy de la variation d'intonation
    - ratio_silence : proportion de fenêtres considérées comme silencieuses (< -35 dB)
    """
    if not prosody_raw:
        return resume
    ratio_silence = prosody_raw.get('ratio_silence', 0)
    db_moyen = prosody_raw.get('db_moyen', -20)
    db_variation = prosody_raw.get('db_variation', 0)
    resume['ratio_silence_pct'] = round(ratio_silence * 100, 1)
    resume['db_moyen'] = db_moyen
    resume['db_variation'] = db_variation
    resume['commentaire_prosodie'] = _commentaire_prosodie(ratio_silence, db_moyen)
    resume['commentaire_intonation'] = _commentaire_intonation(db_variation)
    return resume


def _commentaire_prosodie(ratio_silence: float, db_moyen: float) -> str:
    if ratio_silence > 0.35:
        return (
            f"Nombreuses pauses ({ratio_silence * 100:.0f}% de silence) "
            "— debit hache ou hesitant."
        )
    if ratio_silence > 0.20:
        return f"Pauses moderees ({ratio_silence * 100:.0f}%) — rythme acceptable."
    return f"Debit fluide, peu de silences ({ratio_silence * 100:.0f}%)."


def _commentaire_intonation(db_variation: float) -> str:
    """Interprète l'écart-type du volume comme proxy d'intonation / expressivité vocale."""
    if db_variation > 12:
        return "Voix tres modulee — forte variation de volume (expressivite elevee)."
    if db_variation > 7:
        return "Bonne variation de volume — voix vivante et engagee."
    if db_variation > 3:
        return "Variation de volume moderee — voix correcte mais peu expressive."
    return "Volume quasi constant — voix monotone, peu de variation d'intonation detectee."


def _commentaire_posture(ratio_contact: float) -> str:
    if ratio_contact > 0.7:
        return 'Tres bon contact visuel avec la camera.'
    if ratio_contact > 0.4:
        return 'Contact visuel correct, a ameliorer par moments.'
    return "L'etudiant regarde rarement la camera : posture a travailler."


# ============================================================================
# Extraction de frames depuis une vidéo de démonstration (opencv)
# ============================================================================

def extraire_frames_video(chemin_video: str | Path, n_frames: int = 12) -> list[str]:
    """Extrait n_frames images clés d'une vidéo, encodées en base64 JPEG.
    Essaie imageio+ffmpeg en priorité, puis opencv comme fallback.
    Retourne [] si aucun backend n'est disponible ou si la vidéo est illisible.
    """
    import base64
    import io as _io

    chemin = str(chemin_video)
    frames_b64: list[str] = []

    # ── Méthode 1 : imageio + imageio-ffmpeg (recommandé) ──────────────────
    try:
        import imageio
        from PIL import Image

        reader = imageio.get_reader(chemin, 'ffmpeg')
        meta   = reader.get_meta_data()

        # Estimer le nombre total de frames
        n_total = meta.get('nframes', 0)
        fps      = meta.get('fps', 25) or 25
        duration = meta.get('duration', 0)
        if not n_total or n_total > 1_000_000:
            n_total = int(duration * fps) if duration > 0 else 0

        if n_total > 0:
            indices = [min(int(i * n_total / n_frames), n_total - 1)
                       for i in range(n_frames)]
        else:
            # Durée inconnue : lire les premières frames disponibles
            indices = list(range(n_frames))

        for idx in indices:
            try:
                frame_np = reader.get_data(idx)       # numpy array RGB uint8
                img = Image.fromarray(frame_np)
                # Redimensionner à 768px max (économie tokens)
                w, h = img.size
                if max(w, h) > 768:
                    scale = 768 / max(w, h)
                    img = img.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
                buf = _io.BytesIO()
                img.save(buf, format='JPEG', quality=75)
                frames_b64.append(base64.b64encode(buf.getvalue()).decode('utf-8'))
            except Exception:
                continue

        reader.close()
        if frames_b64:
            logger.info('[video-frames] imageio: %d frames extraites', len(frames_b64))
            return frames_b64

    except ImportError:
        logger.debug('[video-frames] imageio/Pillow non disponible — essai opencv')
    except Exception as exc:
        logger.warning('[video-frames] imageio echec (%s) — essai opencv', exc)

    # ── Méthode 2 : opencv-python-headless (fallback) ──────────────────────
    try:
        import cv2

        cap = cv2.VideoCapture(chemin)
        if not cap.isOpened():
            logger.warning('[video-frames] opencv: impossible d\'ouvrir %s', chemin)
            return []

        total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        if total <= 0:
            cap.release()
            return []

        step = max(1, total // n_frames)
        for i in range(n_frames):
            pos = min(i * step, total - 1)
            cap.set(cv2.CAP_PROP_POS_FRAMES, pos)
            ret, frame = cap.read()
            if not ret:
                continue
            h, w = frame.shape[:2]
            if max(h, w) > 768:
                scale = 768 / max(h, w)
                frame = cv2.resize(frame, (int(w * scale), int(h * scale)))
            ok, buf = cv2.imencode('.jpg', frame, [cv2.IMWRITE_JPEG_QUALITY, 75])
            if ok:
                frames_b64.append(base64.b64encode(buf).decode('utf-8'))

        cap.release()
        logger.info('[video-frames] opencv: %d frames extraites', len(frames_b64))
        return frames_b64

    except ImportError:
        logger.warning('[video-frames] ni imageio ni opencv disponibles — analyse video desactivee')
    except Exception as exc:
        logger.warning('[video-frames] opencv echec: %s', exc)

    return []


# ============================================================================
# Récupération d'un dépôt GitHub public (PyGitHub)
# ============================================================================

def fetch_github_repo(url: str) -> dict:
    """Récupère les métadonnées, README et structure d'un repo GitHub public.
    Retourne un dict exploitable par l'agent Claude.
    """
    try:
        from github import Github, GithubException
    except ImportError:
        logger.warning('[github] PyGitHub non installé')
        return {'erreur': 'PyGitHub non installé'}

    # Normaliser l'URL : https://github.com/owner/repo[/...] → owner/repo
    try:
        chemin = url.rstrip('/').split('github.com/')[-1]
        # Garder seulement owner/repo (ignorer les sous-chemins /blob /tree etc.)
        parties = chemin.split('/')
        repo_path = '/'.join(parties[:2])
    except Exception:
        return {'erreur': f'URL GitHub invalide : {url}'}

    g = Github()  # API publique — 60 req/h sans token
    try:
        repo = g.get_repo(repo_path)
    except GithubException as exc:
        logger.warning('[github] Repo introuvable : %s (%s)', repo_path, exc)
        return {'erreur': f'Repo introuvable ou privé : {repo_path}'}
    except Exception as exc:
        return {'erreur': str(exc)}

    # ── Langages utilisés ─────────────────────────────────────────
    try:
        langages = repo.get_languages()  # dict {lang: nb_bytes}
        total_bytes = sum(langages.values()) or 1
        langages_pct = {
            lang: round(100 * nb / total_bytes, 1)
            for lang, nb in sorted(langages.items(), key=lambda x: -x[1])[:8]
        }
    except Exception:
        langages_pct = {}

    # ── README ────────────────────────────────────────────────────
    readme_texte = ''
    try:
        readme = repo.get_readme()
        import base64 as _b64
        readme_texte = _b64.b64decode(readme.content).decode('utf-8', errors='replace')
        readme_texte = readme_texte[:4000]  # Limiter à 4000 chars
    except Exception:
        readme_texte = '(README absent ou inaccessible)'

    # ── Structure de premier niveau ───────────────────────────────
    try:
        contents = repo.get_contents('')
        structure = [
            f"{'[D]' if c.type == 'dir' else '[F]'} {c.name}"
            for c in sorted(contents, key=lambda x: (x.type != 'dir', x.name))
        ][:30]
    except Exception:
        structure = []

    # ── Statistiques commits ──────────────────────────────────────
    try:
        commits = repo.get_commits()
        nb_commits = commits.totalCount
    except Exception:
        nb_commits = 0

    # ── Fichiers source principaux (max 3 fichiers × 2000 chars) ──
    fichiers_principaux = []
    extensions_code = {'.py', '.js', '.ts', '.java', '.cpp', '.c', '.cs', '.go', '.rs', '.php'}
    try:
        for item in contents:
            if item.type == 'file' and Path(item.name).suffix.lower() in extensions_code:
                import base64 as _b64
                content_raw = repo.get_contents(item.path)
                texte = _b64.b64decode(content_raw.content).decode('utf-8', errors='replace')
                fichiers_principaux.append({
                    'nom': item.name,
                    'contenu': texte[:2000],
                })
                if len(fichiers_principaux) >= 3:
                    break
    except Exception:
        pass

    return {
        'nom': repo.name,
        'description': repo.description or '',
        'url': url,
        'langages': langages_pct,
        'nb_commits': nb_commits,
        'nb_etoiles': repo.stargazers_count,
        'readme': readme_texte,
        'structure': structure,
        'fichiers_principaux': fichiers_principaux,
        'date_creation': str(repo.created_at.date()) if repo.created_at else '',
        'date_dernier_push': str(repo.pushed_at.date()) if repo.pushed_at else '',
    }
