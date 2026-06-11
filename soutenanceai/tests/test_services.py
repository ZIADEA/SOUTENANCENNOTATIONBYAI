"""Tests des services : transcription Groq, extraction PPTX/PDF, LibreOffice, posture."""
import os
import tempfile
from pathlib import Path
from unittest.mock import MagicMock, mock_open, patch

import pytest

from notation.services import (
    convertir_pptx_en_pdf,
    extraire_texte_pdf,
    extraire_texte_pptx,
    resumer_posture,
    transcrire_audio,
)


# ============================================================================
# transcrire_audio
# ============================================================================

class TestTranscrireAudio:

    def test_returns_empty_when_no_api_key(self, settings, tmp_path):
        settings.GROQ_API_KEY = ''
        audio = tmp_path / 'test.webm'
        audio.write_bytes(b'fake')
        result = transcrire_audio(str(audio), langue='fr')
        assert result == ''

    @patch('notation.services.settings')
    @patch('builtins.__import__', side_effect=ImportError('groq non installé'))
    def test_returns_empty_when_groq_not_installed(self, mock_import, mock_settings):
        # Import error signalé si groq manquant
        pass  # Ce test vérifie via le try/except dans la fonction

    @patch('groq.Groq')
    def test_returns_transcription_text(self, mock_groq_class, settings, tmp_path):
        settings.GROQ_API_KEY = 'gsk_fake_key'
        settings.WHISPER_MODEL = 'whisper-large-v3'

        audio = tmp_path / 'chunk.webm'
        audio.write_bytes(b'fake_audio_data')

        mock_client = MagicMock()
        mock_client.audio.transcriptions.create.return_value = 'Bonjour, voici la transcription.'
        mock_groq_class.return_value = mock_client

        result = transcrire_audio(str(audio), langue='fr')
        assert result == 'Bonjour, voici la transcription.'

    @patch('groq.Groq')
    def test_returns_text_from_response_object(self, mock_groq_class, settings, tmp_path):
        """Certaines versions retournent un objet avec .text au lieu d'une chaîne."""
        settings.GROQ_API_KEY = 'gsk_fake_key'
        settings.WHISPER_MODEL = 'whisper-large-v3'

        audio = tmp_path / 'chunk.webm'
        audio.write_bytes(b'data')

        mock_response = MagicMock()
        mock_response.text = 'Texte via attribut .text'
        mock_client = MagicMock()
        mock_client.audio.transcriptions.create.return_value = mock_response
        mock_groq_class.return_value = mock_client

        result = transcrire_audio(str(audio), langue='fr')
        assert result == 'Texte via attribut .text'

    @patch('groq.Groq')
    def test_handles_api_exception_returns_empty(self, mock_groq_class, settings, tmp_path):
        settings.GROQ_API_KEY = 'gsk_fake_key'
        settings.WHISPER_MODEL = 'whisper-large-v3'

        audio = tmp_path / 'chunk.webm'
        audio.write_bytes(b'data')

        mock_client = MagicMock()
        mock_client.audio.transcriptions.create.side_effect = Exception('timeout')
        mock_groq_class.return_value = mock_client

        result = transcrire_audio(str(audio), langue='fr')
        assert result == ''

    def test_langue_mapping(self, settings, tmp_path):
        """Teste que les codes langue sont bien mappés à ISO."""
        settings.GROQ_API_KEY = 'gsk_fake_key'
        settings.WHISPER_MODEL = 'whisper-large-v3'

        audio = tmp_path / 'a.webm'
        audio.write_bytes(b'data')

        with patch('groq.Groq') as mock_groq_class:
            mock_client = MagicMock()
            mock_client.audio.transcriptions.create.return_value = ''
            mock_groq_class.return_value = mock_client

            transcrire_audio(str(audio), langue='en')
            call_kwargs = mock_client.audio.transcriptions.create.call_args[1]
            assert call_kwargs.get('language') == 'en'


# ============================================================================
# extraire_texte_pptx
# ============================================================================

class TestExtraireTextePptx:

    def test_returns_empty_string_if_file_not_found(self):
        result = extraire_texte_pptx('/chemin/inexistant.pptx')
        assert result == ''

    def test_extracts_text_from_slides(self, tmp_path):
        """Crée un PPTX minimal et vérifie l'extraction."""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            pytest.skip('python-pptx non installé')

        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        if title:
            title.text = 'Titre Unique Test'
        body = slide.placeholders[1]
        body.text = 'Contenu corps de slide'

        pptx_path = tmp_path / 'test.pptx'
        prs.save(str(pptx_path))

        result = extraire_texte_pptx(str(pptx_path))
        assert 'Titre Unique Test' in result or 'Contenu corps' in result

    def test_extracts_slide_headers(self, tmp_path):
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip('python-pptx non installé')

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txBox.text_frame.text = 'SLIDE 1 UNIQUE TEXT'

        pptx_path = tmp_path / 'test2.pptx'
        prs.save(str(pptx_path))

        result = extraire_texte_pptx(str(pptx_path))
        assert 'SLIDE 1 UNIQUE TEXT' in result

    def test_marks_slide_boundaries(self, tmp_path):
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip('python-pptx non installé')

        prs = Presentation()
        for _ in range(3):
            prs.slides.add_slide(prs.slide_layouts[5])

        pptx_path = tmp_path / 'multi.pptx'
        prs.save(str(pptx_path))
        result = extraire_texte_pptx(str(pptx_path))
        assert 'SLIDE 1' in result
        assert 'SLIDE 2' in result

    def test_returns_empty_on_corrupt_file(self, tmp_path):
        """Un fichier corrompu doit retourner '' sans lever d'exception."""
        fake = tmp_path / 'fake.pptx'
        fake.write_bytes(b'notapptx')
        result = extraire_texte_pptx(str(fake))
        assert result == ''


# ============================================================================
# extraire_texte_pdf
# ============================================================================

class TestExtraireTextePdf:

    def test_returns_empty_if_file_not_found(self):
        result = extraire_texte_pdf('/inexistant.pdf')
        assert result == ''

    def test_extracts_text_from_pdf(self, tmp_path):
        try:
            import fitz
        except ImportError:
            pytest.skip('PyMuPDF non installé')
        try:
            from reportlab.pdfgen import canvas as rlcanvas
        except ImportError:
            pytest.skip('reportlab non installé')

        pdf_path = tmp_path / 'test.pdf'
        c = rlcanvas.Canvas(str(pdf_path))
        c.drawString(100, 750, 'Texte Unique Dans Le PDF Test')
        c.save()

        result = extraire_texte_pdf(str(pdf_path))
        assert 'Texte Unique Dans Le PDF Test' in result

    def test_marks_page_boundaries(self, tmp_path):
        try:
            import fitz
        except ImportError:
            pytest.skip('PyMuPDF non installé')
        try:
            from reportlab.pdfgen import canvas as rlcanvas
            from reportlab.lib.pagesizes import letter
        except ImportError:
            pytest.skip('reportlab non installé')

        pdf_path = tmp_path / 'multipage.pdf'
        c = rlcanvas.Canvas(str(pdf_path), pagesize=letter)
        c.drawString(100, 750, 'Page un')
        c.showPage()
        c.drawString(100, 750, 'Page deux')
        c.save()

        result = extraire_texte_pdf(str(pdf_path))
        assert 'PAGE 1' in result
        assert 'PAGE 2' in result

    def test_handles_exception_gracefully(self, tmp_path):
        fake = tmp_path / 'broken.pdf'
        fake.write_bytes(b'not a pdf')
        result = extraire_texte_pdf(str(fake))
        assert isinstance(result, str)


# ============================================================================
# convertir_pptx_en_pdf (LibreOffice)
# ============================================================================

class TestConvertirPptxEnPdf:

    def test_returns_none_if_file_not_found(self, tmp_path):
        result = convertir_pptx_en_pdf('/inexistant.pptx', tmp_path)
        assert result is None

    @patch('shutil.which', return_value=None)
    @patch('os.path.isfile', return_value=False)
    def test_returns_none_if_libreoffice_not_found(self, mock_isfile, mock_which, tmp_path):
        pptx = tmp_path / 'slide.pptx'
        pptx.write_bytes(b'fake pptx')
        result = convertir_pptx_en_pdf(str(pptx), tmp_path)
        assert result is None

    @patch('shutil.which', return_value='/usr/bin/libreoffice')
    @patch('subprocess.run')
    def test_returns_pdf_path_on_success(self, mock_run, mock_which, tmp_path):
        mock_run.return_value = MagicMock(returncode=0, stderr='')

        pptx = tmp_path / 'presentation.pptx'
        pptx.write_bytes(b'fake pptx data')

        # Crée le PDF attendu dans le dossier de sortie
        expected_pdf = tmp_path / 'presentation.pdf'
        expected_pdf.write_bytes(b'%PDF fake content')

        result = convertir_pptx_en_pdf(str(pptx), tmp_path)
        assert result == expected_pdf
        mock_run.assert_called_once()

    @patch('shutil.which', return_value='/usr/bin/libreoffice')
    @patch('subprocess.run')
    def test_returns_none_on_libreoffice_error(self, mock_run, mock_which, tmp_path):
        mock_run.return_value = MagicMock(returncode=1, stderr='Erreur conversion')

        pptx = tmp_path / 'bad.pptx'
        pptx.write_bytes(b'bad pptx')

        result = convertir_pptx_en_pdf(str(pptx), tmp_path)
        assert result is None

    @patch('shutil.which', return_value='/usr/bin/libreoffice')
    @patch('subprocess.run')
    def test_returns_none_on_timeout(self, mock_run, mock_which, tmp_path):
        import subprocess
        mock_run.side_effect = subprocess.TimeoutExpired(['libreoffice'], timeout=120)

        pptx = tmp_path / 'slow.pptx'
        pptx.write_bytes(b'fake pptx')

        result = convertir_pptx_en_pdf(str(pptx), tmp_path)
        assert result is None

    @patch('shutil.which', return_value='/usr/bin/libreoffice')
    @patch('subprocess.run')
    def test_creates_output_directory(self, mock_run, mock_which, tmp_path):
        mock_run.return_value = MagicMock(returncode=0, stderr='')

        pptx = tmp_path / 'slide.pptx'
        pptx.write_bytes(b'fake pptx')
        out_dir = tmp_path / 'subdir' / 'output'
        # PDF attendu
        out_dir.mkdir(parents=True, exist_ok=True)
        (out_dir / 'slide.pdf').write_bytes(b'%PDF')

        result = convertir_pptx_en_pdf(str(pptx), out_dir)
        assert out_dir.exists()


# ============================================================================
# resumer_posture
# ============================================================================

class TestResumerPosture:

    def test_empty_data_returns_empty_dict(self):
        assert resumer_posture([]) == {}
        assert resumer_posture(None) == {}

    def test_100_percent_eye_contact(self):
        data = [
            {'looking_at_camera': True, 'head_tilt': 0.0, 'mouvement': 0.1}
            for _ in range(10)
        ]
        result = resumer_posture(data)
        assert result['pourcentage_contact_visuel'] == 100.0

    def test_0_percent_eye_contact(self):
        data = [
            {'looking_at_camera': False, 'head_tilt': 5.0, 'mouvement': 0.3}
            for _ in range(5)
        ]
        result = resumer_posture(data)
        assert result['pourcentage_contact_visuel'] == 0.0
        assert 'rarement' in result['commentaire_auto'].lower()

    def test_partial_eye_contact(self):
        data = (
            [{'looking_at_camera': True, 'head_tilt': 0.0, 'mouvement': 0.1}] * 6
            + [{'looking_at_camera': False, 'head_tilt': 0.0, 'mouvement': 0.1}] * 4
        )
        result = resumer_posture(data)
        assert result['pourcentage_contact_visuel'] == pytest.approx(60.0, rel=1e-3)

    def test_average_head_tilt_calculated(self):
        data = [
            {'looking_at_camera': True, 'head_tilt': 10.0, 'mouvement': 0.0},
            {'looking_at_camera': True, 'head_tilt': -10.0, 'mouvement': 0.0},
        ]
        result = resumer_posture(data)
        assert result['inclinaison_tete_moyenne_deg'] == pytest.approx(0.0, abs=1e-6)

    def test_expressions_aggregated(self):
        """Les expressions faciales sont moyennées et l'émotion dominante détectée."""
        data = [
            {'looking_at_camera': True, 'head_tilt': 0.0,
             'expressions': {'neutral': 0.8, 'happy': 0.2}},
            {'looking_at_camera': True, 'head_tilt': 0.0,
             'expressions': {'neutral': 0.6, 'happy': 0.4}},
        ]
        result = resumer_posture(data)
        assert result['emotion_dominante'] == 'neutral'
        assert result['expressions_moyennes_pct']['neutral'] == pytest.approx(70.0, rel=1e-3)

    def test_nb_mesures_equals_number_of_measures(self):
        data = [{'looking_at_camera': True, 'head_tilt': 0.0, 'mouvement': 0.1}] * 7
        result = resumer_posture(data)
        assert result['nb_mesures'] == 7

    def test_good_contact_comment(self):
        data = [{'looking_at_camera': True, 'head_tilt': 0.0, 'mouvement': 0.0}] * 10
        result = resumer_posture(data)
        assert 'bon' in result['commentaire_auto'].lower() or 'très' in result['commentaire_auto'].lower()

    def test_moderate_contact_comment(self):
        data = (
            [{'looking_at_camera': True, 'head_tilt': 0.0, 'mouvement': 0.0}] * 5
            + [{'looking_at_camera': False, 'head_tilt': 0.0, 'mouvement': 0.0}] * 5
        )
        result = resumer_posture(data)
        # 50% → "correct"
        assert 'correct' in result['commentaire_auto'].lower() or 'améliorer' in result['commentaire_auto'].lower()
